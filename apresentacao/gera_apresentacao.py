# -*- coding: utf-8 -*-
"""
Divine Chocolates — Modelo de Acompanhamento de Projetos (P&D)
Gera um .pptx 16:9 totalmente editável:
  1. Capa
  2. Exemplo preenchido — Slide "Acompanhamento do Projeto"
  3. Exemplo preenchido — Slide "Imagens & Orçamento"
  4. Modelo em branco  — Slide "Acompanhamento do Projeto"
  5. Modelo em branco  — Slide "Imagens & Orçamento"
  6. Como usar + legenda
"""
import os
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE as MSO_LINE

# ---------- paleta Divine ----------
CHOC   = RGBColor(0x3E, 0x24, 0x19)   # marrom chocolate escuro
CHOC2  = RGBColor(0x5A, 0x3A, 0x2A)
CHOC3  = RGBColor(0x7A, 0x52, 0x38)
RED_LOGO = RGBColor(0xA8, 0x32, 0x24) # quadradinho da logo
GOLD   = RGBColor(0xC9, 0x93, 0x2E)
GOLDLT = RGBColor(0xE0, 0xA5, 0x3A)
CREAM  = RGBColor(0xFB, 0xF6, 0xEA)
CREAM2 = RGBColor(0xF4, 0xEA, 0xD7)
CREAM3 = RGBColor(0xF9, 0xF3, 0xE4)
LINE   = RGBColor(0xE3, 0xD6, 0xBF)
INK    = RGBColor(0x2A, 0x1B, 0x12)
MUTE   = RGBColor(0x9C, 0x8C, 0x72)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CREAMTXT = RGBColor(0xE7, 0xD8, 0xC2)

G  = RGBColor(0x43, 0xA0, 0x47); GD = RGBColor(0x2E, 0x7D, 0x32); GB = RGBColor(0xE6, 0xF1, 0xE3)
A  = RGBColor(0xE0, 0xA5, 0x3A); AD = RGBColor(0x9A, 0x6A, 0x12); AB = RGBColor(0xFB, 0xEF, 0xD3)
R  = RGBColor(0xC0, 0x39, 0x2F); RD = RGBColor(0xB2, 0x3A, 0x2A); RB = RGBColor(0xF7, 0xE0, 0xDA)
N  = RGBColor(0xCD, 0xBD, 0xA4); ND = RGBColor(0x8A, 0x7A, 0x66)

FONT   = "Segoe UI"
SCRIPT = "Segoe Script"

LOGO_WM = os.path.join(os.path.dirname(os.path.abspath(__file__)), "logo_wordmark.png")
HAS_LOGO = os.path.exists(LOGO_WM)
HEADER_BG = os.path.join(os.path.dirname(os.path.abspath(__file__)), "header_bg.jpg")
HAS_HEADER = os.path.exists(HEADER_BG)

import anim_lib
def _ids(s): return {sp.shape_id for sp in s.shapes}
def _grp(s, before): return sorted(_ids(s) - before)   # ids criados desde o snapshot

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]

# ---------- helpers ----------
def new_slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s

def shape(s, kind, x, y, w, h, fill=None, line=None, line_w=0.75, shadow_off=True, radius=None, dash=None):
    sp = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
        if dash: sp.line.dash_style = dash
    if shadow_off:
        sp.shadow.inherit = False
    if radius is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    return sp

def rect(s, x, y, w, h, fill, **kw):
    return shape(s, MSO_SHAPE.RECTANGLE, x, y, w, h, fill=fill, **kw)

def card(s, x, y, w, h, fill=WHITE):
    return shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, line=LINE, line_w=1.0, radius=0.055)

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=0, line_spacing=None, wrap=True, autosize_off=True):
    """runs: lista de parágrafos; cada parágrafo é lista de (texto, size, bold, color, italic, font)"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after: p.space_after = Pt(space_after)
        if line_spacing: p.line_spacing = line_spacing
        for (t, size, bold, color, *rest) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = rest[1] if len(rest) > 1 and rest[1] else FONT
            r.font.italic = bool(rest[0]) if rest else False
    return tb

def pill(s, x, y, w, h, fill, label, tcolor, size=11, line_c=None):
    p = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, line=line_c, radius=0.5)
    tf = p.text_frame
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
    r = para.add_run(); r.text = label
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = tcolor; r.font.name = FONT
    return p

def status_dot(s, cx, cy, st, d=0.24):
    """st: 'g' concluído · 'a' em andamento · 'r' bloqueado · 'n' não iniciado"""
    x, y = cx - d/2, cy - d/2
    if st == 'g':
        o = shape(s, MSO_SHAPE.OVAL, x, y, d, d, fill=G)
        sym, sc = "✓", WHITE
    elif st == 'a':
        o = shape(s, MSO_SHAPE.OVAL, x, y, d, d, fill=A)
        sym, sc = None, WHITE
    elif st == 'r':
        o = shape(s, MSO_SHAPE.OVAL, x, y, d, d, fill=R)
        sym, sc = "!", WHITE
    else:
        o = shape(s, MSO_SHAPE.OVAL, x, y, d, d, fill=WHITE, line=N, line_w=1.5)
        sym, sc = None, ND
    if sym:
        tf = o.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = sym
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = sc; r.font.name = FONT
    return o

def logo_block(s, x=0.30, y=0.17, sz=0.62):
    if HAS_LOGO:
        pic = s.shapes.add_picture(LOGO_WM, Inches(x + 0.02), Inches(0.22), height=Inches(0.54))
        return x + 0.02 + pic.width / 914400 + 0.18
    lg = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, sz, sz, fill=RED_LOGO, radius=0.22)
    tf = lg.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Divine"
    r.font.size = Pt(10); r.font.bold = True; r.font.italic = True
    r.font.color.rgb = GOLDLT; r.font.name = SCRIPT
    return 1.08

def header(s, eyebrow, title, prioridade=None, status_geral=None, status_fill=GOLDLT, status_tc=INK):
    if HAS_HEADER:
        s.shapes.add_picture(HEADER_BG, Inches(0), Inches(0), width=Inches(13.333), height=Inches(0.98))
    else:
        rect(s, 0, 0, 13.333, 0.98, CHOC)
    rect(s, 0, 0.98, 13.333, 0.028, GOLD)          # filete dourado
    tx0 = logo_block(s)
    text(s, tx0, 0.155, 8.0, 0.25, [[(eyebrow, 10, True, GOLDLT)]])
    text(s, tx0 - 0.02, 0.375, 8.0, 0.55, [[(title, 25, True, WHITE)]])
    if prioridade is not None:
        text(s, 9.35, 0.14, 1.55, 0.2, [[("PRIORIDADE", 8, True, CREAMTXT)]], align=PP_ALIGN.CENTER)
        pill(s, 9.35, 0.40, 1.55, 0.40, WHITE, prioridade, CHOC, size=11.5)
    if status_geral is not None:
        text(s, 11.10, 0.14, 1.90, 0.2, [[("STATUS GERAL", 8, True, CREAMTXT)]], align=PP_ALIGN.CENTER)
        pill(s, 11.10, 0.40, 1.90, 0.40, status_fill, status_geral, status_tc, size=11.5)

def subband(s, resp, abertura, atualizado, link=""):
    rect(s, 0, 1.008, 13.333, 0.40, CREAM2)
    tem_link = bool((link or "").strip())
    text(s, 0.38, 1.075, 5.8 if tem_link else 8.4, 0.28, [[
        ("Responsável: ", 10, True, CHOC), (resp, 10, False, CHOC2),
        ("      Data de abertura: ", 10, True, CHOC), (abertura, 10, False, CHOC2)]])
    if tem_link:
        url = link if link.lower().startswith("http") else "https://" + link
        chip = pill(s, 6.35, 1.045, 3.35, 0.32, WHITE, "🔎  Pesquisa de mercado  ↗", CHOC, size=9.5, line_c=GOLD)
        try:
            run = chip.text_frame.paragraphs[0].runs[0]
            run.hyperlink.address = url
        except Exception:
            pass
    text(s, 9.9, 1.075, 3.05, 0.28, [[("Atualizado: " + atualizado, 9.5, False, MUTE, True)]],
         align=PP_ALIGN.RIGHT)

def footer(s, legend=True):
    rect(s, 0, 7.18, 13.333, 0.32, CHOC)
    if HAS_LOGO:
        pic = s.shapes.add_picture(LOGO_WM, Inches(12.0), Inches(7.23), height=Inches(0.22))
        pic.left = Inches(12.95) - pic.width
    else:
        text(s, 12.15, 7.215, 0.95, 0.26, [[("Divine", 11, True, GOLDLT, True, SCRIPT)]],
             align=PP_ALIGN.RIGHT, wrap=False)
    if legend:
        items = [(G, "Concluído"), (A, "Em andamento"), (None, "Não iniciado"), (R, "Bloqueado")]
        x = 0.38
        for col, lab in items:
            if col is None:
                shape(s, MSO_SHAPE.OVAL, x, 7.28, 0.12, 0.12, fill=WHITE, line=N, line_w=1.0)
            else:
                shape(s, MSO_SHAPE.OVAL, x, 7.28, 0.12, 0.12, fill=col)
            text(s, x + 0.18, 7.235, 1.5, 0.22, [[(lab, 8.5, True, CREAMTXT)]], wrap=False)
            x += 0.18 + len(lab) * 0.093 + 0.38

def card_header(s, x, y, icon_fill, icon_sym, icon_tc, label, label_color=CHOC):
    ic = shape(s, MSO_SHAPE.OVAL, x, y, 0.26, 0.26, fill=icon_fill)
    tf = ic.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = icon_sym
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = icon_tc; r.font.name = FONT
    text(s, x + 0.37, y + 0.015, 4.5, 0.26, [[(label, 11.5, True, label_color)]])

def bullets(s, x, y, w, h, items, color=CHOC2, size=10.5):
    paras = [[("•  ", size, True, GOLD), (t, size, False, color)] for t in items]
    text(s, x, y, w, h, paras, space_after=5, line_spacing=1.12)

# ============================================================
# SLIDE 1 — CAPA
# ============================================================
s = new_slide(CHOC)
CAPA_BG = os.path.join(os.path.dirname(os.path.abspath(__file__)), "capa_bg.jpg")
if os.path.exists(CAPA_BG):
    s.shapes.add_picture(CAPA_BG, Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
else:
    shape(s, MSO_SHAPE.OVAL, 10.4, -1.6, 4.6, 4.6, fill=RGBColor(0x4A, 0x2E, 0x1E))
    shape(s, MSO_SHAPE.OVAL, 11.5, 5.3, 3.4, 3.4, fill=RGBColor(0x47, 0x2B, 0x1C))
rect(s, 0, 7.34, 13.333, 0.16, GOLD)

# capa cinematográfica: cada elemento entra em sequência (fade)
_cg = []
_b = _ids(s)
text(s, 0.95, 1.55, 10.0, 0.35, [[("P & D   ·   P E S Q U I S A   &   D E S E N V O L V I M E N T O", 12, True, GOLDLT)]])
_cg.append(_grp(s, _b)); _b = _ids(s)
text(s, 0.92, 2.0, 11.0, 1.9, [
    [("Acompanhamento", 45, True, WHITE)],
    [("de Projetos", 45, True, WHITE)]], line_spacing=1.0)
_cg.append(_grp(s, _b)); _b = _ids(s)
text(s, 0.95, 4.05, 8.6, 0.75, [[
    ("Status do portfólio de novos produtos — andamento, pendências e próximos passos, projeto a projeto.",
     15, False, CREAMTXT)]], line_spacing=1.25)
_cg.append(_grp(s, _b)); _b = _ids(s)
p = pill(s, 0.95, 5.0, 3.3, 0.52, None, "Reunião  ·  Junho / 2026", RGBColor(0xF4, 0xE6, 0xCC), size=13)
p.line.color.rgb = GOLD; p.line.width = Pt(1.25)
_cg.append(_grp(s, _b)); _b = _ids(s)
if HAS_LOGO:
    _pic = s.shapes.add_picture(LOGO_WM, Inches(0.92), Inches(5.85), height=Inches(1.0))
    text(s, 0.92, 6.95, _pic.width / 914400, 0.3,
         [[("C H O C O L A T E   D E   V E R D A D E", 9.5, True, CREAMTXT)]],
         align=PP_ALIGN.CENTER, wrap=False)
else:
    text(s, 0.92, 6.05, 4.0, 0.85, [[("Divine", 38, True, GOLDLT, True, SCRIPT)]], wrap=False)
    text(s, 0.98, 6.92, 5.0, 0.3, [[("C H O C O L A T E   D E   V E R D A D E", 9.5, True, CREAMTXT)]], wrap=False)
_cg.append(_grp(s, _b))
anim_lib.set_entrance(s, _cg)

# ============================================================
# builders dos 2 slides de projeto
# ============================================================
# nomes SEM número — a numeração é automática (segue a ordem)
ETAPAS_PADRAO = ["Pesquisa de mercado", "Desenv. da formulação", "Custo e preço de venda",
                 "Viabilidade produtiva/teste em linha", "Homologação", "Volume (comercial)",
                 "Engenharias e FTs", "Criação da arte e etiquetas", "Compras", "Tabela comercial",
                 "Produção", "Comunicado vendedores", "Campanha de lançamento"]

def slide_acompanhamento(nome, prioridade, status_geral, status_fill, status_tc,
                         resp, abertura, atualizado, etapas, pendencias, acoes, observ, link=""):
    s = new_slide()
    header(s, "ACOMPANHAMENTO DO PROJETO", nome, prioridade, status_geral, status_fill, status_tc)
    subband(s, resp, abertura, atualizado, link)

    # ---- card esquerdo: status detalhado ----
    card(s, 0.35, 1.62, 6.55, 5.42)
    text(s, 0.62, 1.82, 4.0, 0.26, [[("STATUS DETALHADO", 11.5, True, CHOC)]])

    tx, ty, tw = 0.62, 2.18, 6.02
    col_w = [2.85, 0.72, 1.45, 1.00]
    # altura/fonte adaptativas: cabem até ~13 etapas sem estourar o card
    hdr_h = 0.32
    row_h = min(0.442, max(0.30, 4.42 / max(len(etapas), 1)))
    efs = 8.5 if len(etapas) > 11 else 9.5
    tbl = s.shapes.add_table(len(etapas) + 1, 4, Inches(tx), Inches(ty),
                             Inches(tw), Inches(hdr_h + row_h * len(etapas))).table
    tbl.first_row = False; tbl.horz_banding = False
    for c, w in enumerate(col_w): tbl.columns[c].width = Inches(w)
    tbl.rows[0].height = Inches(hdr_h)
    for i in range(len(etapas)): tbl.rows[i + 1].height = Inches(row_h)

    def cell(rr, cc, t, size, bold, color, align=PP_ALIGN.LEFT, fill=None):
        c = tbl.cell(rr, cc)
        c.fill.solid(); c.fill.fore_color.rgb = fill if fill else WHITE
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.07); c.margin_right = Inches(0.03)
        c.margin_top = c.margin_bottom = Inches(0.01)
        pgh = c.text_frame.paragraphs[0]; pgh.alignment = align
        r = pgh.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT

    for cidx, htxt in enumerate(["ETAPA", "STATUS", "RESPONSÁVEL", "PREVISTO"]):
        cell(0, cidx, htxt, 8, True, CHOC3,
             PP_ALIGN.CENTER if cidx in (1, 3) else PP_ALIGN.LEFT, CREAM2)
    for i, (etapa, st, resp_e, prev) in enumerate(etapas):
        f = CREAM3 if i % 2 else WHITE
        cell(i + 1, 0, f"{i + 1}. {etapa}", efs, True, CHOC, PP_ALIGN.LEFT, f)
        cell(i + 1, 1, "", 9, False, CHOC, PP_ALIGN.CENTER, f)
        cell(i + 1, 2, resp_e, efs, False, CHOC2, PP_ALIGN.LEFT, f)
        cell(i + 1, 3, prev, efs, False, CHOC2, PP_ALIGN.CENTER, f)
        status_dot(s, tx + col_w[0] + col_w[1] / 2, ty + hdr_h + row_h * i + row_h / 2, st, d=0.235)

    # ---- coluna direita (entra em cascata: Pendências -> Ações -> Observações) ----
    _ag = []
    _b = _ids(s)
    card(s, 7.10, 1.62, 5.85, 1.66)
    card_header(s, 7.36, 1.83, A, "!", WHITE, "PENDÊNCIAS", AD)
    bullets(s, 7.46, 2.24, 5.28, 0.95, pendencias)
    _ag.append(_grp(s, _b)); _b = _ids(s)

    card(s, 7.10, 3.42, 5.85, 1.66)
    card_header(s, 7.36, 3.63, G, "→", WHITE, "PRÓXIMAS AÇÕES", GD)
    bullets(s, 7.46, 4.04, 5.28, 0.95, acoes)
    _ag.append(_grp(s, _b)); _b = _ids(s)

    card(s, 7.10, 5.22, 5.85, 1.82)
    card_header(s, 7.36, 5.43, GOLDLT, "✎", WHITE, "OBSERVAÇÕES E COMENTÁRIOS", CHOC)
    text(s, 7.46, 5.85, 5.28, 1.05, [[(observ, 10.5, False, CHOC2, True)]], line_spacing=1.2)
    _ag.append(_grp(s, _b))
    anim_lib.set_entrance(s, _ag)

    footer(s)

    # ---- paleta de carimbos (fora do slide: aparece só na edição, nunca na apresentação) ----
    card(s, 13.55, 1.62, 1.55, 3.1, fill=WHITE)
    text(s, 13.67, 1.74, 1.35, 0.5, [[("PALETA", 9, True, CHOC)], [("copie e cole ↓", 8, False, MUTE, True)]])
    for i, (st, lab) in enumerate([('g', "Concluído"), ('a', "Andamento"), ('r', "Bloqueado"), ('n', "Não inic.")]):
        yy = 2.35 + i * 0.56
        status_dot(s, 13.85, yy + 0.12, st, d=0.235)
        text(s, 14.05, yy + 0.02, 1.0, 0.22, [[(lab, 8, False, CHOC2)]], wrap=False)
    return s

def slide_imagens_orcamento(nome, orc):
    s = new_slide()
    header(s, "IMAGENS DO PRODUTO & ORÇAMENTO", nome)

    # ---- área de imagem (tracejada) — grupo 1 ----
    _og = []; _b = _ids(s)
    ph = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.35, 1.30, 7.55, 5.60,
               fill=WHITE, line=GOLDLT, line_w=1.5, radius=0.03, dash=MSO_LINE.DASH)
    cx = 0.35 + 7.55 / 2
    ic = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx - 0.42, 3.05, 0.84, 0.84, fill=GOLDLT, radius=0.25)
    shape(s, MSO_SHAPE.OVAL, cx - 0.24, 3.20, 0.17, 0.17, fill=WHITE)
    shape(s, MSO_SHAPE.ISOSCELES_TRIANGLE, cx - 0.26, 3.42, 0.52, 0.36, fill=WHITE)
    text(s, 0.85, 4.10, 6.55, 0.32, [[("Espaço para imagem do produto", 14, True, CHOC)]],
         align=PP_ALIGN.CENTER)
    text(s, 0.85, 4.46, 6.55, 0.3, [[
        ("(insira aqui fotos, mockups, protótipos ou capturas de tela)", 10, False, MUTE, True)]],
        align=PP_ALIGN.CENTER)

    _og.append(_grp(s, _b)); _b = _ids(s)
    # ---- card orçamento — grupo 2 ----
    card(s, 8.10, 1.30, 4.90, 5.60)
    card_header(s, 8.36, 1.52, GOLD, "$", WHITE, "ORÇAMENTO", CHOC)

    fx, fw = 8.40, 4.30
    fy = 1.98
    for lab, val in orc["campos"]:
        text(s, fx, fy, fw, 0.18, [[(lab, 8, True, MUTE)]])
        fld = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, fx, fy + 0.21, fw, 0.34,
                    fill=CREAM2, line=LINE, line_w=0.75, radius=0.16)
        tf = fld.text_frame
        tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.05)
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pgh = tf.paragraphs[0]; pgh.alignment = PP_ALIGN.LEFT
        r = pgh.add_run(); r.text = val
        r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = CHOC2; r.font.name = FONT
        fy += 0.72

    text(s, fx, fy, fw, 0.18, [[("OBSERVAÇÕES SOBRE O ORÇAMENTO", 8, True, MUTE)]])
    obs = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, fx, fy + 0.21, fw, 1.12,
                fill=RGBColor(0xFD, 0xF3, 0xDC), line=RGBColor(0xE8, 0xCF, 0x93), line_w=0.75, radius=0.06)
    rect(s, fx, fy + 0.21, 0.055, 1.12, GOLDLT)
    tf = obs.text_frame
    tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.07); tf.margin_bottom = Inches(0.05)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    pgh = tf.paragraphs[0]; pgh.alignment = PP_ALIGN.LEFT
    r = pgh.add_run(); r.text = orc["obs"]
    r.font.size = Pt(10); r.font.color.rgb = CHOC2; r.font.name = FONT
    _og.append(_grp(s, _b))
    anim_lib.set_entrance(s, _og)

    footer(s, legend=False)
    return s

# ============================================================
# SLIDES 2–3 — EXEMPLO PREENCHIDO (Barrinhas FOR FIT)
# ============================================================
_st_forfit = ['g', 'a', 'a', 'r', 'g', 'n', 'a', 'g', 'n', 'a', 'n', 'n', 'n']
_rp_forfit = ["Trade", "P&D", "Camila", "Carla e Samuel", "CQ", "Mateus", "P&D",
              "Marketing", "—", "P&D", "—", "Marketing", "Marketing"]
_pv_forfit = ["08/06", "24/07", "24/07", "03/08", "17/07", "—", "24/07",
              "24/07", "—", "17/07", "—", "—", "—"]
etapas_forfit = [(ETAPAS_PADRAO[i], _st_forfit[i], _rp_forfit[i], _pv_forfit[i])
                 for i in range(len(ETAPAS_PADRAO))]
slide_acompanhamento(
    "Barrinhas FOR FIT", "ALTA", "EM TESTE", AB, AD,
    "Paloma e Isabele", "15/06/2026", date.today().strftime("%d/%m/%Y"),
    etapas_forfit,
    ["Aguardar fechamento de formulação;", "Moinho de rolos;"],
    ["Concluir ajustes no moinho de rolos;", "Fechar formulações;"],
    "Projeto dentro do cronograma.",
    link="www.exemplo.com/pesquisa-barrinhas")

slide_imagens_orcamento("Barrinhas FOR FIT", {
    "campos": [("CUSTO ESTIMADO", "R$ —"), ("CUSTO REAL", "R$ —"),
               ("FORNECEDOR(ES)", "—"), ("APROVAÇÃO", "—")],
    "obs": "A definir junto ao fechamento de formulação e custo (etapas 2 e 3)."})

# ============================================================
# SLIDES 4–5 — MODELO EM BRANCO
# ============================================================
etapas_blank = [(e, 'n', "—", "—") for e in ETAPAS_PADRAO]
slide_acompanhamento(
    "[ Nome do Projeto ]", "—", "PLANEJAMENTO", CREAM2, CHOC,
    "—", "—", "—",
    etapas_blank,
    ["—", "—"],
    ["—", "—"],
    "—")

slide_imagens_orcamento("[ Nome do Projeto ]", {
    "campos": [("CUSTO ESTIMADO", "R$ —"), ("CUSTO REAL", "R$ —"),
               ("FORNECEDOR(ES)", "—"), ("APROVAÇÃO", "—")],
    "obs": "—"})

# ============================================================
# SLIDE 6 — COMO USAR
# ============================================================
s = new_slide()
header(s, "GUIA RÁPIDO", "Como usar este modelo")

card(s, 0.35, 1.30, 7.55, 5.60)
text(s, 0.62, 1.52, 5.0, 0.26, [[("PASSO A PASSO", 11.5, True, CHOC)]])
passos = [
    ("1.", "Novo projeto:", " selecione os dois slides-modelo “[ Nome do Projeto ]”, duplique (Ctrl+D) e mova o par para o final."),
    ("2.", "Cabeçalho:", " troque o nome do projeto, a prioridade e o status geral (a caixinha dourada é só texto + cor de preenchimento)."),
    ("3.", "Status das etapas:", " clique na bolinha → guia Formatar → Preenchimento da Forma → escolha a cor (verde, âmbar, vermelho ou branco). Para o “✓”, copie uma bolinha verde já pronta."),
    ("4.", "Tabela:", " edite etapa, responsável e previsão como texto normal; linhas podem ser adicionadas ou removidas."),
    ("5.", "Slide de imagens:", " arraste a foto para dentro da área tracejada (ou Inserir → Imagens) e apague o ícone de espaço reservado."),
    ("6.", "Orçamento:", " preencha os campos de custo, fornecedor e aprovação; use a caixa amarela para observações."),
]
paras = [[(n + "  ", 11, True, GOLD), (b, 11, True, CHOC), (t, 11, False, CHOC2)] for n, b, t in passos]
text(s, 0.62, 1.95, 7.0, 4.8, paras, space_after=10, line_spacing=1.18)

card(s, 8.10, 1.30, 4.90, 5.60)
text(s, 8.36, 1.52, 4.2, 0.26, [[("LEGENDA DE STATUS", 11.5, True, CHOC)]])
leg = [('g', "Concluído", "Etapa finalizada e validada."),
       ('a', "Em andamento", "Etapa em execução."),
       ('r', "Bloqueado", "Travado — depende de decisão ou terceiro."),
       ('n', "Não iniciado", "Ainda não começou.")]
ly = 2.05
for st, lab, desc in leg:
    status_dot(s, 8.55, ly + 0.14, st, d=0.28)
    text(s, 8.85, ly, 4.0, 0.24, [[(lab, 11.5, True, CHOC)]])
    text(s, 8.85, ly + 0.26, 4.0, 0.24, [[(desc, 9.5, False, MUTE)]])
    ly += 0.72

text(s, 8.36, 5.05, 4.2, 0.26, [[("STATUS GERAL (cabeçalho)", 11.5, True, CHOC)]])
gerais = [("EM ANDAMENTO", GOLDLT, INK), ("CONCLUÍDO", G, WHITE), ("PAUSADO", N, INK),
          ("EM TESTE", AB, AD), ("PLANEJAMENTO", CREAM2, CHOC)]
gx, gy = 8.36, 5.42
for lab, f, tc in gerais:
    w = 0.42 + len(lab) * 0.085
    if gx + w > 12.85:
        gx = 8.36; gy += 0.52
    pill(s, gx, gy, w, 0.38, f, lab, tc, size=9.5)
    gx += w + 0.16
footer(s, legend=False)

# ============================================================
# SLIDE FINAL — MUITO OBRIGADA
# ============================================================
s = new_slide(CHOC)
if os.path.exists(CAPA_BG):
    s.shapes.add_picture(CAPA_BG, Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
else:
    shape(s, MSO_SHAPE.OVAL, 10.4, -1.6, 4.6, 4.6, fill=RGBColor(0x4A, 0x2E, 0x1E))
    shape(s, MSO_SHAPE.OVAL, 11.5, 5.3, 3.4, 3.4, fill=RGBColor(0x47, 0x2B, 0x1C))
rect(s, 0, 7.34, 13.333, 0.16, GOLD)
_tg = []; _b = _ids(s)
text(s, 0.92, 2.55, 11.5, 1.3, [[("Muito obrigada!", 54, True, WHITE)]])
_tg.append(_grp(s, _b)); _b = _ids(s)
text(s, 0.95, 3.95, 10.0, 0.4, [[("Divine Chocolates · Pesquisa & Desenvolvimento", 15, False, CREAMTXT)]])
_tg.append(_grp(s, _b)); _b = _ids(s)
if HAS_LOGO:
    _pic = s.shapes.add_picture(LOGO_WM, Inches(0.92), Inches(5.55), height=Inches(0.9))
    text(s, 0.92, 6.52, _pic.width / 914400, 0.3,
         [[("C H O C O L A T E   D E   V E R D A D E", 9.5, True, CREAMTXT)]],
         align=PP_ALIGN.CENTER, wrap=False)
_tg.append(_grp(s, _b))
anim_lib.set_entrance(s, _tg)

# transição fade suave em todos os slides
for _sl in prs.slides:
    anim_lib.add_fade_transition(_sl)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Divine_Acompanhamento_Projetos_MODELO.pptx")
prs.save(out)
print("OK:", out)
